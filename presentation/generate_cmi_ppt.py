# -*- coding: utf-8 -*-
"""
CMI Motor Architecture Exploration - PPT generator.

발표 원고(표지 + 본문 13장 + Appendix)를 편집 가능한 .pptx 로 생성한다.
- 형식      : python-pptx 로 생성한 편집 가능한 PowerPoint
- 언어      : 영문 제목 + 한글 본문 (원문 그대로)
- 시각자료  : 실제 사진/차트는 라벨이 달린 플레이스홀더 박스로 자리만 확보
              (추후 presentation/images/ 의 이미지를 add_picture 로 교체)
- 디자인    : 회사 템플릿으로 옮기기 쉬운 단순·중립 레이아웃

실행:  python presentation/generate_cmi_ppt.py
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:  # 점선 테두리(있으면 사용, 없으면 실선)
    from pptx.enum.line import MSO_LINE_DASH_STYLE
    _HAS_DASH = True
except Exception:  # pragma: no cover
    _HAS_DASH = False


# --------------------------------------------------------------------------- #
# 스타일 상수 (중립 톤 — 회사 템플릿 적용 전 베이스)
# --------------------------------------------------------------------------- #
KOR_FONT = "맑은 고딕"          # 한글 본문
ENG_FONT = "Calibri"            # 영문 제목

DARK = RGBColor(0x23, 0x2B, 0x38)      # 제목/강조 텍스트
TEXT = RGBColor(0x3A, 0x3F, 0x47)      # 본문 텍스트
MUTED = RGBColor(0x8A, 0x92, 0x9C)     # 보조 텍스트
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)    # 액센트(차분한 블루)
ACCENT_DK = RGBColor(0x16, 0x4E, 0x7E) # 진한 액센트
KEY_BG = RGBColor(0xEA, 0xF2, 0xFA)    # Key Message 박스 배경
PH_FILL = RGBColor(0xF1, 0xF3, 0xF6)   # 플레이스홀더 배경
PH_LINE = RGBColor(0xB4, 0xBC, 0xC6)   # 플레이스홀더 테두리
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 슬라이드(16:9)
SW, SH = 13.333, 7.5


# --------------------------------------------------------------------------- #
# 저수준 헬퍼
# --------------------------------------------------------------------------- #
def _set_run_font(run, name, size, bold=False, color=None, italic=False):
    """run 의 latin/ea/cs typeface 를 모두 같은 폰트로 지정 → 한글이 맑은 고딕으로 렌더."""
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = name
    if color is not None:
        f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_textbox(slide, left, top, width, height, anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb, tf


def fill_tf(tf, paras):
    """paras: list of dict(text, size, bold, color, name, align, space_after, space_before, level)."""
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if pa.get("space_after") is not None:
            p.space_after = Pt(pa["space_after"])
        if pa.get("space_before") is not None:
            p.space_before = Pt(pa["space_before"])
        p.level = pa.get("level", 0)
        run = p.add_run()
        run.text = pa["text"]
        _set_run_font(
            run,
            pa.get("name", KOR_FONT),
            pa["size"],
            pa.get("bold", False),
            pa.get("color", TEXT),
            pa.get("italic", False),
        )


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE, dashed=False):
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dashed and _HAS_DASH:
            try:
                shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            except Exception:
                pass
    return shp


def add_placeholder(slide, left, top, width, height, label):
    """이미지/차트가 들어갈 자리를 표시하는 라벨 박스."""
    shp = add_rect(slide, left, top, width, height, fill=PH_FILL, line=PH_LINE,
                   line_w=1.0, dashed=True)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    paras = [{"text": "［ 이미지 영역 ］", "size": 9.5, "bold": False,
              "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 3}]
    for ln in str(label).split("\n"):
        paras.append({"text": ln, "size": 11, "bold": True, "color": RGBColor(0x5A, 0x62, 0x6C),
                      "align": PP_ALIGN.CENTER, "space_after": 1})
    fill_tf(tf, paras)
    return shp


def slide_title(slide, title_text):
    """상단 영문 타이틀 + 액센트 라인."""
    _, tf = add_textbox(slide, 0.55, 0.38, 12.23, 0.9, anchor=MSO_ANCHOR.MIDDLE)
    fill_tf(tf, [{"text": title_text, "size": 26, "bold": True, "color": DARK,
                  "name": ENG_FONT, "space_after": 0}])
    add_rect(slide, 0.57, 1.28, 4.2, 0.045, fill=ACCENT)
    add_rect(slide, 4.77, 1.28, 8.0, 0.045, fill=RGBColor(0xDD, 0xE2, 0xE8))


def future_chip(slide, left, top, width, text):
    shp = add_rect(slide, left, top, width, 0.62, fill=ACCENT,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.12)
    fill_tf(tf, [{"text": "▶  Future Opportunity", "size": 9.5, "bold": True,
                  "color": RGBColor(0xCF, 0xE3, 0xF5), "space_after": 1},
                 {"text": text, "size": 12.5, "bold": True, "color": WHITE, "space_after": 0}])
    return shp


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --------------------------------------------------------------------------- #
# 슬라이드 빌더
# --------------------------------------------------------------------------- #
def add_title_slide(prs):
    s = blank_slide(prs)
    add_rect(s, 0.0, 0.0, SW, 0.32, fill=ACCENT)
    add_rect(s, 0.0, SH - 0.22, SW, 0.22, fill=RGBColor(0xE6, 0xEA, 0xEF))

    _, tf = add_textbox(s, 1.0, 2.55, 11.33, 1.5, anchor=MSO_ANCHOR.BOTTOM)
    fill_tf(tf, [{"text": "CMI Motor Architecture Exploration", "size": 40, "bold": True,
                  "color": DARK, "name": ENG_FONT, "align": PP_ALIGN.LEFT, "space_after": 0}])

    add_rect(s, 1.03, 4.18, 2.6, 0.06, fill=ACCENT)

    _, tf2 = add_textbox(s, 1.0, 4.35, 11.33, 1.2)
    fill_tf(tf2, [
        {"text": "CMI 기반 차세대 Motor Architecture 탐색", "size": 18, "bold": True,
         "color": ACCENT_DK, "space_after": 4},
        {"text": "기존 소재로 구현하기 어려운 모터 구조를 가능하게 하는 자성 플랫폼 제안",
         "size": 13, "bold": False, "color": MUTED, "space_after": 0},
    ])
    return s


def add_content_slide(prs, spec):
    s = blank_slide(prs)
    title = "{}. {}".format(spec["num"], spec["title"])
    slide_title(s, title)

    # ----- 좌측 칼럼 : Key Message + Supporting Data -----
    # Key Message 강조 박스
    kbox = add_rect(s, 0.55, 1.5, 5.78, 1.45, fill=KEY_BG, line=None,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ktf = kbox.text_frame
    ktf.word_wrap = True
    ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ktf.margin_left = Inches(0.16)
    ktf.margin_right = Inches(0.16)
    ktf.margin_top = Inches(0.08)
    ktf.margin_bottom = Inches(0.08)
    fill_tf(ktf, [
        {"text": "KEY MESSAGE", "size": 10, "bold": True, "color": ACCENT, "space_after": 3},
        {"text": spec["key"], "size": 14.5, "bold": True, "color": DARK, "space_after": 0},
    ])

    # Supporting Data
    _, dtf = add_textbox(s, 0.6, 3.12, 5.73, 3.05)
    paras = [{"text": spec["data_label"], "size": 11.5, "bold": True, "color": ACCENT,
              "space_after": 6}]
    for b in spec["bullets"]:
        paras.append({"text": "•  " + b, "size": 13, "bold": False, "color": TEXT,
                      "space_after": 5})
    fill_tf(dtf, paras)

    # Future Opportunity chip
    future_chip(s, 0.55, 6.36, 5.78, spec["future"])

    # ----- 우측 칼럼 : Visual Evidence -----
    _, vtf = add_textbox(s, 6.78, 1.5, 6.0, 0.4)
    fill_tf(vtf, [{"text": "VISUAL EVIDENCE", "size": 10, "bold": True, "color": MUTED,
                   "space_after": 0}])

    labels = spec["visuals"]
    n = len(labels)
    area_top, area_bottom = 1.98, 7.02
    gap = 0.16
    box_h = (area_bottom - area_top - gap * (n - 1)) / n
    for i, lab in enumerate(labels):
        top = area_top + i * (box_h + gap)
        add_placeholder(s, 6.78, top, 6.0, box_h, lab)
    return s


def add_closing_slide(prs, spec):
    s = blank_slide(prs)
    slide_title(s, "{}. {}".format(spec["num"], spec["title"]))

    # Key Message (full width)
    kbox = add_rect(s, 0.55, 1.5, 12.23, 0.95, fill=KEY_BG, line=None,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ktf = kbox.text_frame
    ktf.word_wrap = True
    ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ktf.margin_left = Inches(0.18)
    fill_tf(ktf, [
        {"text": "KEY MESSAGE", "size": 10, "bold": True, "color": ACCENT, "space_after": 2},
        {"text": spec["key"], "size": 16, "bold": True, "color": DARK, "space_after": 0},
    ])

    # 두 칼럼
    def column(left, heading, items):
        _, tf = add_textbox(s, left, 2.75, 5.9, 2.35)
        paras = [{"text": heading, "size": 13, "bold": True, "color": ACCENT_DK,
                  "space_after": 7}]
        for it in items:
            paras.append({"text": "•  " + it, "size": 13.5, "bold": False, "color": TEXT,
                          "space_after": 5})
        fill_tf(tf, paras)

    column(0.6, "Seeking Partners", spec["partners"])
    column(6.85, "Expected Collaboration", spec["collab"])

    # Final Message 박스
    fbox = add_rect(s, 0.55, 5.35, 12.23, 1.6, fill=ACCENT_DK, line=None,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ftf = fbox.text_frame
    ftf.word_wrap = True
    ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ftf.margin_left = Inches(0.22)
    ftf.margin_right = Inches(0.22)
    fill_tf(ftf, [
        {"text": "FINAL MESSAGE", "size": 10, "bold": True,
         "color": RGBColor(0xBF, 0xD7, 0xEC), "space_after": 5},
        {"text": spec["final"][0], "size": 14, "bold": True, "color": WHITE, "space_after": 4},
        {"text": spec["final"][1], "size": 14, "bold": True, "color": WHITE, "space_after": 0},
    ])
    return s


def add_appendix_slide(prs, items, note):
    s = blank_slide(prs)
    slide_title(s, "Appendix — 준비 항목")

    _, sub = add_textbox(s, 0.6, 1.45, 12.2, 0.4)
    fill_tf(sub, [{"text": "발표 후 보강 예정 데이터 (Data to Prepare)", "size": 13,
                   "bold": True, "color": MUTED, "space_after": 0}])

    half = (len(items) + 1) // 2
    cols = [items[:half], items[half:]]
    for ci, col in enumerate(cols):
        _, tf = add_textbox(s, 0.7 + ci * 6.15, 2.05, 5.95, 3.6)
        paras = []
        for it in col:
            paras.append({"text": "□  " + it, "size": 14, "bold": False, "color": TEXT,
                          "space_after": 8})
        fill_tf(tf, paras)

    nbox = add_rect(s, 0.55, 5.85, 12.23, 1.0, fill=KEY_BG, line=None,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ntf = nbox.text_frame
    ntf.word_wrap = True
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ntf.margin_left = Inches(0.2)
    fill_tf(ntf, [
        {"text": "Motor Prototype 결과", "size": 11, "bold": True, "color": ACCENT,
         "space_after": 3},
        {"text": note, "size": 13.5, "bold": True, "color": DARK, "space_after": 0},
    ])
    return s


# --------------------------------------------------------------------------- #
# 콘텐츠 데이터
# --------------------------------------------------------------------------- #
CONTENT_SLIDES = [
    {
        "num": "1", "title": "Why Are We Looking at Motors?",
        "key": "CMI는 인덕터용 기술로 개발되었지만, 개발 과정에서 모터 분야 적용 가능성을 발견하였다.",
        "data_label": "현재 확보 실적",
        "bullets": ["50kW CMI", "100kW CMI", "3상 일체형 자기회로", "냉각 구조 통합 경험"],
        "visuals": ["100kW 3상 일체형 CMI 시제품 / 단품 코아 사진",
                     "50kW CMI 시제품 / 단품 코아 사진",
                     "자기회로 일체형 · 몰딩형 구조 사진"],
        "future": "Motor Magnetic Architecture",
    },
    {
        "num": "2", "title": "What We Learned from Power Magnetics",
        "key": "CMI는 기존 자성체와 다른 구조적 장점을 제공한다.",
        "data_label": "현재 확보 데이터",
        "bullets": ["Density 5.3~6.0 g/cc", "Thermal Conductivity 3~4 W/mK",
                     "Bmax 12,000~15,000 Gauss", "Operating Temperature -20~155℃",
                     "200℃까지 투자율 유지 특성"],
        "visuals": ["소재 특성표", "온도별 투자율 변화 그래프", "신뢰성 시험 결과"],
        "future": "고온 환경 모터 적용",
    },
    {
        "num": "3", "title": "Why Conventional Motors Are Difficult to Evolve",
        "key": "현재 모터 산업의 과제는 소재보다 구조 혁신에 있다.",
        "data_label": "시장 트렌드",
        "bullets": ["Axial Flux Motor", "Integrated Drive Unit", "Robot Joint Motor",
                     "Compact Servo"],
        "visuals": ["Axial Motor 단면", "Hairpin Motor", "Robot Joint Motor"],
        "future": "구조 혁신 기반 모터 개발",
    },
    {
        "num": "4", "title": "Existing Magnetic Material Landscape",
        "key": "CMI는 Electrical Steel 대체재가 아니다.",
        "data_label": "Material Positioning  (비투자율 μr)",
        "bullets": ["Electrical Steel  —  μr : 3,000~6,000", "SMC  —  μr : 200~800",
                     "CMI  —  μr : 20~60", "Hard Magnet  —  NdFeB"],
        "visuals": ["Material Position Map\n(X축: Shape Freedom  /  Y축: Magnetic Performance)"],
        "future": "새로운 소재 조합 — 규소강·벌크 하이브리드화로 투자율·Bs 상승 가능",
    },
    {
        "num": "5", "title": "Then Why CMI?",
        "key": "CMI의 경쟁력은 투자율이 아니라 형상 자유도이다.",
        "data_label": "CMI 특징",
        "bullets": ["Net Shape", "Integrated Structure", "Embedded Cooling",
                     "Multi-Material Compatibility"],
        "visuals": ["실제 CMI 구조", "100kW 시제품", "3상 일체형 자기회로"],
        "future": "기존 방식으로 구현 불가능한 구조",
    },
    {
        "num": "6", "title": "What Makes CMI Different?",
        "key": "CMI는 자기회로와 구조물을 동시에 설계할 수 있다.",
        "data_label": "비교",
        "bullets": ["Steel  —  높은 투자율 / 낮은 형상 자유도",
                     "CMI  —  낮은 투자율 / 높은 형상 자유도"],
        "visuals": ["Steel Stator", "SMC Stator", "CMI Concept"],
        "future": "Architecture Driven Design",
    },
    {
        "num": "7", "title": "New Possibility #1 – Integrated Coil + Magnetic Circuit",
        "key": "CMI는 권선과 자기회로를 동시에 설계할 수 있다.",
        "data_label": "현재 확보 실적",
        "bullets": ["권선 일체형 성형", "대형 권선 구조 경험"],
        "visuals": ["인덕터 구조", "권선 포함 성형 구조", "미래 모터 구조 Concept"],
        "future": "Manufacturing Simplification",
    },
    {
        "num": "8", "title": "New Possibility #2 – Cooling Integrated Motor",
        "key": "CMI는 냉각 구조와 자기회로를 동시에 설계할 수 있다.",
        "data_label": "100kW CMI 데이터",
        "bullets": ["Ambient 50℃", "Forced Air 0.074 kg/s", "Maximum Temperature 83℃"],
        "visuals": ["열 시뮬레이션 결과", "온도 분포", "냉각 경로"],
        "future": "High Power Density Motor",
    },
    {
        "num": "9", "title": "New Possibility #3 – Hard Magnet Integration",
        "key": "CMI는 Hard Magnet과 일체형 제작이 가능하다.",
        "data_label": "현재 진행 현황",
        "bullets": ["국내 모터 스타트업 공동개발", "자석 일체형 구조 검토 중"],
        "visuals": ["Magnet Insert 구조", "Overmolding 구조", "Rotor Concept"],
        "future": "Integrated Rotor Architecture",
    },
    {
        "num": "10", "title": "New Possibility #4 – Axial Flux Motor",
        "key": "CMI는 Axial Motor 분야에서 가장 높은 가능성을 가진다.",
        "data_label": "Axial Motor 요구사항",
        "bullets": ["복잡 형상", "냉각", "일체형 구조", "소형화"],
        "visuals": ["YASA 구조", "Magnax 구조", "CMI 적용 예상 구조"],
        "future": "EV / Aerospace / Robotics",
    },
    {
        "num": "11", "title": "Proposed Hybrid Magnetic Architecture",
        "key": "CMI는 단독 사용보다 Hybrid Structure에서 가치가 크다.",
        "data_label": "가능한 구조",
        "bullets": ["CMI + Hard Magnet", "CMI + FeSi Ball", "CMI + SMC", "CMI + Steel"],
        "visuals": ["각 구조의 Concept Diagram"],
        "future": "Next Generation Magnetic Architecture",
    },
    {
        "num": "12", "title": "Technical Challenges",
        "key": "현재 기술적 한계를 인정하고, 이를 해결하기 위한 개발을 진행 중이다.",
        "data_label": "현재 과제",
        "bullets": ["낮은 투자율", "낮은 Bs", "모터 효율 검증 필요", "토크밀도 검증 필요",
                     "대량생산 검증 필요"],
        "visuals": ["현재 수준 vs 목표 수준 — Radar Chart"],
        "future": "차세대 플랫폼으로 진화",
    },
]

CLOSING_SLIDE = {
    "num": "13", "title": "What We Are Looking For",
    "key": "모터 전문 기업과 공동 검증을 희망한다.",
    "partners": ["Axial Flux Motor", "Robotics Motor", "Servo Motor", "Special Motor",
                 "Integrated Drive Unit"],
    "collab": ["Motor Simulation", "Prototype Development", "Performance Validation",
               "Manufacturing Validation"],
    "final": [
        "CMI는 모터 소재 혁신을 목표로 하지 않는다.",
        "CMI는 기존 소재로는 구현하기 어려운 새로운 Motor Architecture를 가능하게 하는 플랫폼이 되고자 한다.",
    ],
}

APPENDIX_ITEMS = [
    "BH Curve 데이터", "Core Loss 데이터", "Density 데이터", "Thermal Conductivity 데이터",
    "Temperature Stability 데이터", "Reliability 데이터", "SEM 분석", "단면 분석",
    "자석 접합 강도", "Maxwell / JMAG Simulation",
]
APPENDIX_NOTE = "비욘드 로봇(Beyond Robotics) 데이터 입수 예정"


# --------------------------------------------------------------------------- #
# 빌드 & 검증
# --------------------------------------------------------------------------- #
def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    expected_titles = ["CMI Motor Architecture Exploration"]
    add_title_slide(prs)

    for spec in CONTENT_SLIDES:
        add_content_slide(prs, spec)
        expected_titles.append("{}. {}".format(spec["num"], spec["title"]))

    add_closing_slide(prs, CLOSING_SLIDE)
    expected_titles.append("{}. {}".format(CLOSING_SLIDE["num"], CLOSING_SLIDE["title"]))

    add_appendix_slide(prs, APPENDIX_ITEMS, APPENDIX_NOTE)
    expected_titles.append("Appendix — 준비 항목")

    return prs, expected_titles


def main():
    try:  # Windows 콘솔(cp949)에서도 en-dash/한글이 깨지지 않도록
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

    base = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(base, "CMI_Motor_Architecture_Exploration.pptx")

    prs, expected_titles = build()
    prs.save(out)

    # --- 검증: 다시 열어 슬라이드 수 / 제목 확인 ---
    check = Presentation(out)
    n = len(list(check.slides))
    print("저장 완료:", out)
    print("슬라이드 수:", n, "(기대값 15)")
    for i, slide in enumerate(check.slides, 1):
        # 각 슬라이드 첫 텍스트(제목)를 추출
        title_txt = ""
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                title_txt = shp.text_frame.text.strip().split("\n")[0]
                break
        print("  {:>2}. {}".format(i, title_txt))

    assert n == 15, "슬라이드 수가 15가 아닙니다: {}".format(n)
    print("\nOK - 15 slides generated.")


if __name__ == "__main__":
    main()
